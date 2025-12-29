import os
import fitz
import docx
import openpyxl
from pptx import Presentation
from pathlib import Path
import logging

from app.config import Config
from pathlib import Path

logger = logging.getLogger(__name__)

SUPPORTED_EXTENTIONS = Config.ALLOWED_EXTENSIONS

def extract_text_from_pdf(file_path: str) -> str:
    doc= fitz.open(file_path)
    text=""
    for page in doc:
        text+=page.get_text()
    return text

def extract_text_from_docx(file_path: str) -> str:
    doc = docx.Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text+= shape.text +"\n"
        return text

def extract_text_from_xlsx(file_path: str) -> str:
    wb=openpyxl.load_workbook(file_path, data_only=True)
    text=""
    for sheet in wb.worksheets:
        for row in sheet.iter_rows():
            row_text = " ".join([str(cell.value) if cell.value else ""])
            text+=row_text+ "\n"
    return text

def extract_text_from_txt(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as file:
            text= file.read()
        return text

    except FileNotFoundError as fnfe:
        logger.error(f"File not found: {file_path}")
        raise fnfe
    except Exception as e:
        logger.error(f"An error occured: {e}")
        raise e 

def extract_text_from_file(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    try:
        if ext == ".pdf":
            return extract_text_from_pdf(file_path)
        elif ext == ".docx":
            return extract_text_from_docx(file_path)
        elif ext == ".pptx":
            return extract_text_from_pptx(file_path)
        elif ext == ".xlsx":
            return extract_text_from_xlsx(file_path)
        elif ext == ".txt":
            return extract_text_from_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")
    except Exception as e:
        logger.exception(f"Error processing file: {file_path}")
        raise e 
               
                   